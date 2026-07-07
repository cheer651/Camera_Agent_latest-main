from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt


SOURCE = Path(r"C:\Users\chens\Desktop\camera_project\docs\opening_defense_camera_agent_local.pptx")
TARGET = Path(r"C:\Users\chens\Desktop\camera_project\docs\opening_defense_camera_agent_local_en.pptx")


EXACT_RAW = {
    "研究背景": "Background",
    "系统方案": "Architecture",
    "智能体设计": "Agent Design",
    "阶段进展": "Progress",
    "后续计划": "Future Work",
    "智能视频监控 Agent 系统": "Intelligent Video Surveillance Agent System",
    "答辩内容目录": "Agenda",
    "01 研究背景与问题定义": "01 Background and Problem Definition",
    "传统视频全量上云成本高  ｜  人工巡检难以处理长时间、多摄像头视频  ｜  需要本地化、可交互、可追溯的智能安防方案": "Sending full video streams to the cloud is expensive  |  Manual inspection cannot efficiently handle long-duration multi-camera videos  |  A local, interactive, and traceable intelligent security solution is needed",
    "02 系统方案": "02 System Architecture",
    "两路 RTSP 摄像头接入研究院内网  ｜  本地笔记本完成录制、抽帧、分析与存储  ｜  Web 端统一展示关键帧、问答、日报和告警": "Two RTSP cameras are connected to the institute LAN  |  The local laptop completes recording, frame extraction, analysis, and storage  |  The web UI uniformly displays keyframes, QA results, daily reports, and alerts",
    "03 安防智能体设计": "03 Security Agent Design",
    "LangGraph 多节点工作流  ｜  感知-记忆-推理-执行闭环  ｜  多模型分工协作与结构化输出约束": "LangGraph multi-node workflow  |  Perception-memory-reasoning-action closed loop  |  Multi-model collaboration with structured output constraints",
    "04 当前阶段进展与验证": "04 Current Progress and Validation",
    "已完成两路摄像头、本地 UI、飞书、SQLite/Qdrant  ｜  已接入 Qwen2.5-VL、Qwen2.5、Qwen3-Embedding  ｜  评估当前效果与尚待优化问题": "Two cameras, local UI, Feishu, and SQLite/Qdrant have been implemented  |  Qwen2.5-VL, Qwen2.5, and Qwen3-Embedding have been integrated  |  The current effect and remaining issues are being evaluated",
    "05 后续语音模块与知识库计划": "05 Voice Module and Knowledge Base Plan",
    "大厅麦克风唤醒安防智能体  ｜  构建人员特征本地知识库  ｜  实现‘今天我有没有来过’类自然问答": "Wake up the security agent with a lobby microphone  |  Build a local knowledge base of personnel characteristics  |  Support natural questions such as 'Have I come here today?'",
    "研究背景与问题定义": "Background and Problem Definition",
    "痛点 1：全量视频直接上云成本高": "Pain Point 1: Sending full video streams to the cloud is expensive",
    "多路摄像头持续产生日志与视频流": "Multiple cameras continuously generate logs and video streams",
    "带宽、云端存储、大量静态画面造成重复传输和重复计算": "Bandwidth and cloud storage are costly, and a large number of static frames lead to redundant transmission and repeated computation",
    "痛点 2：事后检索依赖人工回看": "Pain Point 2: Post-event retrieval still depends on manual video review",
    "长时间录像只能按时间线手工拖拽筛查": "Long recordings can only be reviewed manually by dragging along the timeline",
    "难以按‘黑衣人员/几人路过/高危事件’语义查询": "It is difficult to query semantically for 'a person in black / how many people passed by / high-risk events'",
    "缺少面向安保值班人员的自然语言入口": "There is no natural-language interface tailored to on-duty security staff",
    "痛点 3：告警与知识沉淀不足": "Pain Point 3: Alerting and knowledge accumulation are insufficient",
    "异常发现、告警推送、日报总结之间未形成闭环": "Anomaly detection, alert delivery, and daily reporting have not yet formed a closed loop",
    "缺少可长期积累的本地人员/事件知识库": "A local personnel and event knowledge base for long-term accumulation is still missing",
    "难以把‘看见了什么’转为可复用的业务记忆": "It is difficult to convert 'what the system saw' into reusable operational memory",
    "拟解决的核心问题": "Core Research Questions",
    "1. 如何在本地设备上对两路摄像头视频做低冗余关键帧提取，避免把大量无效静态画面直接交给大模型？": "1. How can low-redundancy keyframe extraction be performed on two camera streams locally so that a large number of ineffective static frames are not directly sent to the large model?",
    "2. 如何用多模态智能体将关键帧分析、结构化入库、语义检索、自然语言问答、飞书告警串成可维护的闭环工作流？": "2. How can a multimodal agent connect keyframe analysis, structured storage, semantic retrieval, natural-language QA, and Feishu alerts into a maintainable closed-loop workflow?",
    "3. 如何在后续接入语音模块与人员知识库后，让系统支持‘今天谁来过/我有没有出现过’这类更自然的大厅交互？": "3. After adding a voice module and a personnel knowledge base, how can the system support more natural lobby interactions such as 'Who came here today?' and 'Did I appear today?'",
    "系统架构": "System Architecture",
    "TP-Link 网络摄像头": "TP-Link IP Cameras",
    "RTSP 子码流": "RTSP Sub-stream",
    "研究院内网直连": "Directly Connected within the Institute LAN",
    "视频录制层": "Video Recording Layer",
    "按片段录制 MP4": "Record MP4 Clips in Segments",
    "本地缓存原始视频": "Cache Raw Video Locally",
    "支持 Web 实时预览": "Support Real-time Preview on the Web UI",
    "关键帧提取层": "Keyframe Extraction Layer",
    "OpenCV 运动检测": "OpenCV Motion Detection",
    "窗口内峰值帧选择": "Peak-frame Selection within a Sliding Window",
    "重复帧过滤与事件合并": "Duplicate-frame Filtering and Event Merging",
    "多模态分析层": "Multimodal Analysis Layer",
    "Qwen2.5-VL 输出结构化 JSON": "Qwen2.5-VL Outputs Structured JSON",
    "Pydantic 字段校验": "Pydantic Field Validation",
    "风险等级与事件描述生成": "Risk Grading and Event Description Generation",
    "记忆与交互层": "Memory and Interaction Layer",
    "自然语言问答": "Natural-Language Question Answering",
    "全天总结与飞书告警": "Daily Summaries and Feishu Alerts",
    "四个已落地业务能力": "Four Implemented Business Functions",
    "关键帧分析报告：按日期、摄像头、时间轴查看关键帧与大模型描述。": "Keyframe analysis reports: view keyframes and model-generated descriptions by date, camera, and timeline.",
    "高危异常飞书报警：高风险事件触发后推送图片与文字。": "Feishu alerts for high-risk anomalies: push images and text when high-risk events are triggered.",
    "全天日志 AI 总结：按早晨/下午/晚上/凌晨四段生成日报。": "AI daily log summary: generate daily reports for morning, afternoon, evening, and late-night periods.",
    "自然语言监控对话：基于 SQLite + Qdrant 检索历史事件并用本地 Qwen2.5 组织回答。": "Natural-language monitoring dialogue: retrieve historical events with SQLite + Qdrant and use local Qwen2.5 to organize answers.",
    "安防智能体：感知-记忆-推理-执行闭环": "Security Agent: Perception-Memory-Reasoning-Action Closed Loop",
    "感知模块": "Perception Module",
    "读取 RTSP/视频片段": "Read RTSP Streams / Video Clips",
    "OpenCV 抽取候选关键帧": "Use OpenCV to Extract Candidate Keyframes",
    "过滤静态冗余画面": "Filter Redundant Static Frames",
    "记忆模块": "Memory Module",
    "SQLite 保存结构化事件": "Store Structured Events in SQLite",
    "Qdrant 存语义向量": "Store Semantic Vectors in Qdrant",
    "为 RAG 与日报提供上下文": "Provide Context for RAG and Daily Reports",
    "推理模块": "Reasoning Module",
    "Qwen2.5-VL 解析画面语义": "Qwen2.5-VL Interprets Scene Semantics",
    "Qwen2.5 生成问答/总结": "Qwen2.5 Generates QA Answers / Summaries",
    "Pydantic 约束输出格式": "Pydantic Constrains the Output Format",
    "执行模块": "Action Module",
    "高危事件飞书告警": "Feishu Alerts for High-risk Events",
    "普通事件入库归档": "Archive Normal Events into the Database",
    "日报生成并回写记忆库": "Generate Daily Reports and Write Back to Memory",
    "为什么这一层要突出“智能体”而不是普通脚本？": "Why Emphasize an 'Agent' Instead of Ordinary Scripts?",
    "状态可追踪：每一帧/每一条用户问题都以 State 在节点间传递，方便调试、回溯和扩展。": "Traceable state: each frame and each user question is passed between nodes as State, making debugging, backtracking, and extension easier.",
    "分支可控制：高危事件走“告警分支”，普通事件走“入库分支”，模型失败走“模型未加载成功”降级分支。": "Controllable branching: high-risk events take the alert branch, normal events take the storage branch, and model failures take the fallback branch labeled 'model failed to load'.",
    "记忆可累积：事件写入 SQLite/Qdrant 后，后续问答、日报、语音 Q&A 都能复用同一份本地知识资产。": "Accumulated memory: once events are written into SQLite/Qdrant, subsequent QA, daily reports, and voice Q&A can reuse the same local knowledge asset.",
    "LangGraph 多模型工作流设计": "LangGraph Multi-model Workflow Design",
    "多个模型分别负责图像理解、文本生成、向量检索，LangGraph 负责节点路由、状态传递和条件分支": "Different models handle image understanding, text generation, and vector retrieval, while LangGraph handles node routing, state passing, and conditional branching.",
    "工作流 A：关键帧分析图": "Workflow A: Keyframe Analysis",
    "入队关键帧": "Enqueue Keyframes",
    "构建中文视觉提示词": "Build Visual Prompt Templates",
    "调用 Qwen2.5-VL": "Call Qwen2.5-VL",
    "解析/校验 JSON": "Parse / Validate JSON",
    "事件合并 + 风险判断 + 写库/告警": "Event Merging + Risk Assessment + Storage / Alerting",
    "工作流 B：自然语言监控对话图": "Workflow B: Natural-language Monitoring Dialogue",
    "用户问题 + 历史会话": "User Question + Chat History",
    "问题改写/时间与摄像头约束解析": "Query Rewriting / Time and Camera Constraint Parsing",
    "SQLite 精确筛选 + Qdrant 语义排序": "Exact Filtering with SQLite + Semantic Ranking with Qdrant",
    "事件匹配/统计聚合": "Event Matching / Statistical Aggregation",
    "Qwen2.5 生成可读回答 + 返回关键帧": "Qwen2.5 Generates Readable Answers + Returns Keyframes",
    "记忆与 RAG：本地结构化库 + 向量库双存储": "Memory and RAG: Dual Local Storage with Structured DB + Vector DB",
    "SQLite 结构化记忆": "SQLite Structured Memory",
    "tasks：一次分析任务的起止时间、触发类型、任务状态。": "tasks: start/end time, trigger type, and task status for each analysis run.",
    "camera_runs：每路摄像头录制出的片段路径、帧数、FPS、状态。": "camera_runs: clip path, frame count, FPS, and status for each camera run.",
    "events：关键帧事件时间、风险等级、人数、动作、衣着颜色、图片/视频路径。": "events: event time, risk level, person count, action, clothing color, and image/video paths for each keyframe event.",
    "summaries：按日期保存全天总结及飞书推送状态。": "summaries: store daily summaries and Feishu delivery status by date.",
    "Qdrant 语义向量记忆": "Qdrant Semantic Vector Memory",
    "把事件描述、风险类型、摄像头、时间上下文编码为向量。": "Encode event descriptions, risk types, cameras, and temporal context as vectors.",
    "支持“黑衣服的人”“有人徘徊吗”这类语义模糊查询。": "Support fuzzy semantic queries such as 'a person in black' or 'was someone loitering?'",
    "与 SQLite 精确过滤组合，实现“先硬约束、再语义排序”。": "Combine with exact filtering in SQLite to achieve 'hard constraints first, semantic ranking second'.",
    "后续人员知识库也可沿用同一套向量检索接口。": "The future personnel knowledge base can reuse the same vector retrieval interface.",
    "Qwen3-Embedding：本地部署，负责日志向量化。": "Qwen3-Embedding: locally deployed for log vectorization.",
    "Qwen2.5：读取检索结果后，按安防汇报口吻生成回答/日报。": "Qwen2.5: reads retrieved results and generates answers / daily reports in a security-reporting style.",
    "Prompt 约束：只基于检索日志回答；无记录则返回“未在监控记录中发现相关异常”。": "Prompt constraint: answer strictly based on retrieved logs; if no record exists, return 'No relevant abnormality was found in the monitoring records.'",
    "目前项目进度": "Current Project Progress",
    "初步验证结果与当前问题": "Preliminary Validation Results and Current Issues",
    "2 路": "2 Streams",
    "内网 RTSP 摄像头接入": "Institute-LAN RTSP Camera Access",
    "4 类": "4 Types",
    "已落地核心业务功能": "Core Business Functions Implemented",
    "3 模型": "3 Models",
    "VL / 文本 / Embedding 协同": "VL / Text / Embedding Collaboration",
    "4 表": "4 Tables",
    "SQLite 事件/任务/录制/日报": "SQLite Events / Tasks / Recording / Reports",
    "已验证的闭环能力": "Verified End-to-End Capabilities",
    "网页端可启动任务、查看两路画面、按日期查看关键帧时间轴与分析描述。": "The web UI can start tasks, display two camera views, and show keyframe timelines and analysis results by date.",
    "高危事件可触发飞书推送，普通事件结构化入 SQLite 并写入向量库。": "High-risk events can trigger Feishu alerts, while normal events are written into SQLite and the vector database in structured form.",
    "自然语言对话支持流式输出、历史会话记录和关键帧图片引用。": "Natural-language dialogue supports streaming output, conversation history, and referenced keyframe images.",
    "日报模块可按早晨/下午/晚上/凌晨四个时段总结并生成综合研判文本。": "The daily report module can summarize morning, afternoon, evening, and late-night periods and generate an integrated assessment.",
    "当前暴露的主要问题": "Key Issues Observed at the Current Stage",
    "关键帧提取仍可能出现“无人重复帧偏多”或“某路人员经过漏提取”的情况。": "Keyframe extraction may still produce too many repeated no-person frames or miss pedestrian pass-by events on one camera.",
    "跨帧同一事件合并还需要继续调时间阈值、相似度阈值和代表帧选择策略。": "Same-event merging across frames still requires further tuning of time thresholds, similarity thresholds, and representative-frame selection strategies.",
    "统计型问答虽然已做 SQL 优先，但人物属性字段质量仍依赖 VL 输出稳定性。": "Although statistical QA now prioritizes SQL, the quality of person-attribute fields still depends on the stability of VL outputs.",
    "下一阶段需要把“人是谁/是否是某人再次出现”从事件记忆进一步扩展到人员知识库记忆。": "The next stage needs to extend from event memory to personnel knowledge memory for questions such as 'who is this person' and 'did this person appear again'.",
    "后续扩展方向一：大厅语音唤醒与常见 Q&A": "Future Direction 1: Lobby Voice Wake-up and Common Q&A",
    "在大厅放置高质量麦克风，通过语音唤醒接入安防智能体，实现面向值班/访客的自然语言问答入口": "Place a high-quality microphone in the lobby and wake up the security agent by voice, creating a natural-language entry point for staff and visitors.",
    "大厅麦克风阵列": "Lobby Microphone Array",
    "唤醒词检测": "Wake-word Detection",
    "近场语音采集": "Near-field Speech Capture",
    "降噪/回声抑制": "Noise Reduction / Echo Suppression",
    "ASR + 意图识别": "ASR + Intent Recognition",
    "语音转文字": "Speech-to-Text",
    "时间/摄像头/人物条件抽取": "Time / Camera / Person Constraint Extraction",
    "问答/查询/播报路由": "QA / Query / Response Routing",
    "智能体检索与推理": "Agent Retrieval and Reasoning",
    "访问事件库/人员库": "Access the Event DB / Personnel DB",
    "RAG 检索": "RAG Retrieval",
    "Qwen2.5 生成回答": "Qwen2.5 Generates Answers",
    "语音播报": "Voice Response",
    "TTS 回复": "TTS Reply",
    "前端同步显示": "Synchronized Web Display",
    "必要时附关键帧": "Attach Keyframes When Needed",
    "目标交互样例": "Target Interaction Examples",
    "“今天大厅有没有谁来过？” → 返回时间段、出现次数、对应摄像头和关键帧。": "\"Has anyone come through the lobby today?\" → Return time periods, occurrence counts, related cameras, and keyframes.",
    "“今天我有没有来过？” → 结合人员知识库中的个人描述/特征，检索相似事件后给出结果。": "\"Have I come here today?\" → Combine personal descriptions / features in the personnel knowledge base, retrieve similar events, and return the result.",
    "“今天有没有陌生人多次徘徊？” → 从事件库中按动作类型、人数和时间聚合回答。": "\"Did any unfamiliar person loiter multiple times today?\" → Aggregate answers from the event database by action type, person count, and time.",
    "工程注意点": "Engineering Considerations",
    "麦克风建议优先选带远场拾音、降噪和 USB/网口稳定接入能力的型号。": "The microphone should preferably support far-field pickup, noise reduction, and stable USB / Ethernet access.",
    "语音模块只做入口层，后端仍复用现有 LangGraph 问答工作流和本地知识库。": "The voice module should only serve as the entry layer; the backend should continue reusing the existing LangGraph QA workflow and local knowledge base.",
    "需增加唤醒词误触发控制、权限边界、日志留存和隐私合规说明。": "Additional control is needed for false wake-up triggers, permission boundaries, log retention, and privacy compliance.",
    "后续扩展方向二：构建人员本地知识库": "Future Direction 2: Build a Local Personnel Knowledge Base",
    "把“事件日志”进一步升级为“人员档案 + 历史出现轨迹 + 视觉/文本特征索引”的可问答本地知识库": "Upgrade the 'event log' into a queryable local knowledge base consisting of personnel profiles, historical trajectories, and visual / textual feature indexes.",
    "人员档案表": "Personnel Profile Table",
    "person_id：人员唯一编号": "person_id: unique identifier of the person",
    "name / nickname：姓名或备注名": "name / nickname: person name or alias",
    "role：员工/访客/安保/未知": "role: employee / visitor / security / unknown",
    "appearance_text：身高、发型、衣着偏好、携带物等自然语言描述": "appearance_text: natural-language description of height, hairstyle, clothing preference, carried items, etc.",
    "face/body embedding：可选视觉向量特征": "face / body embedding: optional visual embedding features",
    "authorized_zone：允许出现区域": "authorized_zone: permitted area",
    "事件-人员关联表": "Event-Person Association Table",
    "event_id ↔ person_id：把历史关键帧事件与某个已知人员候选关联起来": "event_id ↔ person_id: associate historical keyframe events with a known personnel candidate",
    "match_score：视觉/文本相似度或规则匹配分数": "match_score: visual / text similarity or rule-matching score",
    "is_manual_verified：是否经过人工确认": "is_manual_verified: whether the match has been manually verified",
    "first_seen_at / last_seen_at：第一次/最近一次出现时间": "first_seen_at / last_seen_at: first / most recent appearance time",
    "trajectory_summary：当天或一周内的出现轨迹摘要": "trajectory_summary: summarized appearance trajectory for the day or the week",
    "面向智能体的查询能力": "Agent-oriented Query Capabilities",
    "“今天我有没有来过？”": "\"Have I come here today?\"",
    "“张三今天最后一次出现在哪个摄像头？”": "\"On which camera did Zhang San appear for the last time today?\"",
    "“这周有没有陌生人连续三天晚上出现？”": "\"Did any unfamiliar person appear for three consecutive evenings this week?\"",
    "“穿黑色外套的人是不是上周二也出现过？”": "\"Did the person in the black coat also appear last Tuesday?\"",
    "本质上是“摄像头事件记忆 + 人员知识库 + 语义检索 + 智能体推理”的组合。": "In essence, this is a combination of camera event memory, a personnel knowledge base, semantic retrieval, and agent reasoning.",
    "后续研究计划": "Future Research Plan",
    "4月": "April",
    "优化关键帧提取与事件合并": "Optimize Keyframe Extraction and Event Merging",
    "复盘无人重复帧问题": "Review the repeated no-person frame issue",
    "调参运动阈值/时间窗/相似度": "Tune motion thresholds / time windows / similarity",
    "提升两路摄像头人员经过召回率": "Improve recall of person pass-by events on both cameras",
    "5月": "May",
    "补强事件字段与统计问答": "Strengthen Event Fields and Statistical QA",
    "稳定输出人数/衣着/动作字段": "Stabilize outputs for person count / clothing / action fields",
    "完善统计型 SQL 查询": "Improve statistical SQL queries",
    "增强“上周几/某日晚上”时间解析": "Improve parsing of expressions such as 'last weekday' / 'that evening'",
    "6月": "June",
    "接入语音唤醒与大厅 Q&A": "Integrate Voice Wake-up and Lobby Q&A",
    "选型麦克风与唤醒链路": "Select microphone hardware and the wake-up pipeline",
    "ASR → 智能体 → TTS 闭环": "ASR → Agent → TTS Closed Loop",
    "前端新增语音状态展示": "Add voice-status display to the front end",
    "7月": "July",
    "建设人员本地知识库": "Build the Local Personnel Knowledge Base",
    "设计人员档案/事件关联表": "Design personnel profile and event-association tables",
    "导入常见人员描述和特征": "Import common personnel descriptions and characteristics",
    "支持“我有没有来过/某人轨迹”查询": "Support queries such as 'Have I been here?' / 'trajectory of a specific person'",
    "8月": "August",
    "系统评测与论文整理": "System Evaluation and Thesis Preparation",
    "长时间稳定性测试": "Long-duration stability tests",
    "准确率/召回率/延迟指标评估": "Evaluate accuracy / recall / latency metrics",
    "整理开题后实验结果与论文结构": "Organize post-proposal experiment results and thesis structure",
    "总结与预期成果": "Summary and Expected Outcomes",
    "已完成基础": "Completed Foundation",
    "本地直连两路摄像头。": "Two cameras have been connected directly within the local LAN.",
    "关键帧提取、VL 分析、结构化入库。": "Keyframe extraction, VL analysis, and structured storage have been completed.",
    "自然语言监控问答、四时段日报总结、飞书告警。": "Natural-language monitoring QA, four-period daily summaries, and Feishu alerts are available.",
    "SQLite + Qdrant 的本地记忆架构已经打通。": "The local memory architecture based on SQLite + Qdrant has been fully connected.",
    "开题后重点突破": "Key Breakthroughs After the Proposal",
    "把 LangGraph 智能体工作流从“可跑”推进到“可解释、可评估、可长期稳定运行”。": "Advance the LangGraph agent workflow from 'it can run' to 'it is explainable, evaluable, and stably operable over the long term'.",
    "把事件库扩展为人员知识库，支持更贴近真实值班场景的自然问答。": "Extend the event database into a personnel knowledge base to support more realistic natural QA in real duty scenarios.",
    "补充语音唤醒模块，让大厅场景可以直接对话使用系统。": "Add a voice wake-up module so that the system can be used directly through spoken dialogue in the lobby scenario.",
    "预期成果形式": "Expected Deliverables",
    "一套面向研究院场景的本地化视频安防智能体原型系统。": "A localized video-security agent prototype system oriented to the institute scenario.",
    "一套可复用的“摄像头事件记忆 + RAG + 语音入口”工程实现方案。": "A reusable engineering solution for 'camera event memory + RAG + voice entry'.",
    "一组围绕关键帧提取、事件问答准确性、告警时效性的实验结果。": "A set of experimental results on keyframe extraction, event-QA accuracy, and alert timeliness.",
    "最终论文与答辩演示材料。": "The final thesis and defense presentation materials.",
    "谢谢各位老师！": "Thank you, professors!",
    "敬请批评指正": "Your comments and suggestions are welcome.",
    "答辩人：陈思        导师：待补充": "Candidate: Chen Si        Supervisor: To be added",
}


SUBSTRING_RAW = [
    ("智能视频监控 Agent 系统", "Intelligent Video Surveillance Agent System"),
    ("研究背景", "Background"),
    ("系统方案", "Architecture"),
    ("智能体设计", "Agent Design"),
    ("阶段进展", "Progress"),
    ("后续计划", "Future Work"),
]


def _repair_mojibake(text: str) -> str:
    try:
        return text.encode("latin-1").decode("utf-8")
    except UnicodeError:
        return text


def _build_exact() -> dict[str, str]:
    repaired = {_repair_mojibake(key): value for key, value in EXACT_RAW.items()}
    # Top navigation already has English subtitles underneath; blank the Chinese layer.
    repaired["研究背景"] = ""
    repaired["系统方案"] = ""
    repaired["智能体设计"] = ""
    repaired["阶段进展"] = ""
    repaired["后续计划"] = ""
    # Refine existing English-only strings in the current deck.
    repaired["AI Agent System for Security Surveil lance Based on Visual Large Language Model"] = (
        "AI Agent System for Security Surveillance Based on Visual Large Language Models"
    )
    repaired["Agent"] = "Agent Design"
    repaired["Plans"] = "Future Work"
    repaired["Thanks！"] = "Thanks!"
    repaired.update(
        {
            "TP-Link 网络摄像头\nRTSP 子码流\n研究院内网直连": (
                "TP-Link IP Cameras\nRTSP Sub-stream\nDirect Institute-LAN Connection"
            ),
            "按片段录制 MP4\n本地缓存原始视频\n支持 Web 实时预览": (
                "Clip-based MP4 Recording\nLocal Raw Video Cache\nSupport Real-time Web Preview"
            ),
            "OpenCV 运动检测\n窗口内峰值帧选择\n重复帧过滤与事件合并": (
                "OpenCV Motion Detection\nPeak-frame Selection within a Window\nDuplicate-frame Filtering and Event Merging"
            ),
            "Qwen2.5-VL 输出结构化 JSON\nPydantic 字段校验\n风险等级与事件描述生成": (
                "Qwen2.5-VL Outputs Structured JSON\nPydantic Field Validation\nRisk Grading and Event Description Generation"
            ),
            "SQLite + Qdrant + Embedding\n自然语言问答\n全天总结与飞书告警": (
                "SQLite + Qdrant + Embedding\nNatural-Language QA\nDaily Summaries and Feishu Alerts"
            ),
            "读取 RTSP/视频片段\nOpenCV 抽取候选关键帧\n过滤静态冗余画面": (
                "Read RTSP Streams / Video Clips\nOpenCV Extracts Candidate Keyframes\nFilter Redundant Static Frames"
            ),
            "SQLite 保存结构化事件\nQdrant 存语义向量\n为 RAG 与日报提供上下文": (
                "SQLite Stores Structured Events\nQdrant Stores Semantic Vectors\nProvide Context for RAG and Daily Reports"
            ),
            "Qwen2.5-VL 解析画面语义\nQwen2.5 生成问答/总结\nPydantic 约束输出格式": (
                "Qwen2.5-VL Interprets Scene Semantics\nQwen2.5 Generates QA Answers / Summaries\nPydantic Constrains Output Format"
            ),
            "高危事件飞书告警\n普通事件入库归档\n日报生成并回写记忆库": (
                "Feishu Alerts for High-risk Events\nArchive Normal Events into the Database\nGenerate Daily Reports and Write Back to Memory"
            ),
            "模块": "Module",
            "当前进度": "Current Status",
            "已实现内容": "Implemented Content",
            "下一步优化": "Next-step Optimization",
            "摄像头接入与录制": "Camera Access and Recording",
            "已完成": "Completed",
            "Cam_01 / Cam_02 RTSP 接入、片段录制、网页预览": (
                "Cam_01 / Cam_02 RTSP access, clip recording, and web preview"
            ),
            "提升掉线重连与长时间运行稳定性": (
                "Improve reconnection after disconnection and long-duration stability"
            ),
            "关键帧提取": "Keyframe Extraction",
            "已完成原型": "Prototype Completed",
            "运动检测、窗口峰值帧、重复帧过滤、事件合并": (
                "Motion detection, peak-frame selection, duplicate-frame filtering, and event merging"
            ),
            "进一步减少无人物重复帧，提高人员经过召回率": (
                "Further reduce repeated no-person frames and improve recall of passing pedestrians"
            ),
            "关键帧分析智能体": "Keyframe Analysis Agent",
            "已接入": "Integrated",
            "Qwen2.5-VL + 中文结构化 JSON + 风险分级": (
                "Qwen2.5-VL + Chinese structured JSON + risk grading"
            ),
            "优化提示词和字段质量，增强人物动作/衣着描述稳定性": (
                "Optimize prompts and field quality, and improve stability of person action / clothing descriptions"
            ),
            "自然语言对话": "Natural-Language Dialogue",
            "已上线": "Online",
            "基于 SQLite + Qdrant 的历史事件检索与流式回答": (
                "Historical event retrieval and streaming answers based on SQLite + Qdrant"
            ),
            "增强复杂时间表达、统计类问答、多人/人物属性推理": (
                "Improve complex time expressions, statistical QA, and multi-person / attribute reasoning"
            ),
            "全天日志总结": "Daily Log Summary",
            "按四时段总结 + 全天综合总结 + 飞书推送": (
                "Four-period summaries + whole-day integrated summary + Feishu push"
            ),
            "强化重点异常提炼和巡检建议生成": (
                "Strengthen key anomaly extraction and patrol-recommendation generation"
            ),
            "前端交互界面": "Front-end Interaction UI",
            "四个 Tag 页面、聊天历史、流式输出、关键帧时间轴": (
                "Four tag pages, chat history, streaming output, and keyframe timelines"
            ),
            "继续修正聊天气泡细节，补充语音入口和人员库页面": (
                "Continue refining chat bubble details and add voice entry and personnel-knowledge-base pages"
            ),
            "唤醒词检测\n近场语音采集\n降噪/回声抑制": (
                "Wake-word Detection\nNear-field Voice Capture\nNoise Reduction / Echo Suppression"
            ),
            "语音转文字\n时间/摄像头/人物条件抽取\n问答/查询/播报路由": (
                "Speech-to-Text\nTime / Camera / Person Constraint Extraction\nQA / Query / Announcement Routing"
            ),
            "访问事件库/人员库\nRAG 检索\nQwen2.5 生成回答": (
                "Access Event DB / Personnel DB\nRAG Retrieval\nQwen2.5 Generates Answers"
            ),
            "TTS 回复\n前端同步显示\n必要时附关键帧": (
                "TTS Reply\nSynchronized Front-end Display\nAttach Keyframes When Needed"
            ),
        }
    )
    return repaired


def _build_substring() -> list[tuple[str, str]]:
    return [(_repair_mojibake(old), new) for old, new in SUBSTRING_RAW]


EXACT = _build_exact()
SUBSTRING = _build_substring()


def translate_text(text: str) -> str:
    stripped = text.strip()
    if not stripped:
        return text
    if stripped in EXACT:
        translated = EXACT[stripped]
    else:
        translated = text
        for old, new in SUBSTRING:
            translated = translated.replace(old, new)
    return translated


def translate_text_frame(text_frame) -> None:
    for para in text_frame.paragraphs:
        if para.text:
            new_text = translate_text(para.text)
            if new_text != para.text:
                para.text = new_text


def _fit_text_frame(text_frame, max_size: float) -> None:
    if not text_frame.text.strip():
        return
    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    text_frame.margin_left = Pt(3)
    text_frame.margin_right = Pt(3)
    text_frame.margin_top = Pt(2)
    text_frame.margin_bottom = Pt(2)
    try:
        text_frame.fit_text(font_family="Arial", max_size=max_size)
    except Exception:
        pass
    for para in text_frame.paragraphs:
        for run in para.runs:
            run.font.name = "Arial"
            run.font.size = Pt(max_size)


def _shape_max_size(slide_index: int, shape, is_table: bool = False, row_idx: int | None = None) -> float:
    if is_table:
        return 9.0 if row_idx == 0 else 8.0
    top = shape.top
    width = shape.width
    height = shape.height
    if slide_index == 0:
        if top < 2600000 and height >= 500000:
            return 20.0
        if top < 3300000:
            return 14.0
        if top < 3800000:
            return 10.0
        return 9.0
    if top < 500000 and height <= 260000:
        return 7.5
    if top < 500000:
        return 13.0
    if top < 1300000 and height <= 320000:
        return 16.0
    if height <= 300000 and width <= 2600000:
        return 10.0
    if height <= 650000 and width <= 2600000:
        return 8.5
    if height <= 700000 and width > 6000000:
        return 10.0
    if height > 2600000:
        return 8.5
    if height > 1500000:
        return 9.5
    return 10.0


def walk_shapes(shapes, slide_index: int) -> None:
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            walk_shapes(shape.shapes, slide_index)
        elif getattr(shape, "has_table", False):
            for row_idx, row in enumerate(shape.table.rows):
                for cell in row.cells:
                    translate_text_frame(cell.text_frame)
                    _fit_text_frame(cell.text_frame, _shape_max_size(slide_index, shape, is_table=True, row_idx=row_idx))
        elif getattr(shape, "has_text_frame", False):
            translate_text_frame(shape.text_frame)
            _fit_text_frame(shape.text_frame, _shape_max_size(slide_index, shape))


def main() -> None:
    prs = Presentation(str(SOURCE))
    for slide_index, slide in enumerate(prs.slides):
        walk_shapes(slide.shapes, slide_index)
    prs.save(str(TARGET))
    print(TARGET)


if __name__ == "__main__":
    main()
